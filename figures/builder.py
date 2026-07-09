from pathlib import Path
import matplotlib.pyplot as plt
import matplotlib.image as mpimg
import pandas as pd

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def collect_figure_files(exp_path):
    exp_path = Path(exp_path)
    fig_root = exp_path / "Figures"
    items = []

    if not fig_root.exists():
        return items

    for p in fig_root.rglob("*.png"):
        if "PublicationFigures" in p.parts:
            continue
        label = str(p.relative_to(fig_root)).replace("\\", "/")
        items.append({"label": label, "path": str(p)})

    return sorted(items, key=lambda x: x["label"])


def suggest_next_figure_name(exp_path):
    exp_path = Path(exp_path)
    pub_root = exp_path / "PublicationFigures"
    pub_root.mkdir(parents=True, exist_ok=True)

    existing = [p.name for p in pub_root.iterdir() if p.is_dir()]
    nums = []

    for name in existing:
        if name.lower().startswith("figure"):
            tail = name[6:]
            try:
                nums.append(int(tail))
            except ValueError:
                pass

    next_num = max(nums) + 1 if nums else 1
    return f"Figure{next_num}"


def generate_caption(panel_labels):
    lines = ["Figure. Composite electrochemical analysis figure."]

    for letter, label in panel_labels:
        lower = label.lower()

        if "raw_overlay" in lower:
            desc = "Raw electrochemical response overlay."
        elif "baseline" in lower:
            desc = "Baseline-corrected electrochemical response."
        elif "delta" in lower or "peak" in lower:
            desc = "Concentration-dependent signal change."
        elif "nyquist" in lower:
            desc = "Nyquist plots obtained from EIS analysis."
        elif "rct" in lower:
            desc = "Charge-transfer resistance response."
        elif "fit" in lower or "langmuir" in lower or "hill" in lower:
            desc = "Model fitting result."
        else:
            desc = "Analysis result."

        lines.append(f"({letter}) {desc}")

    return "\n".join(lines)


def get_layout(n_panels, requested_layout="auto"):
    if requested_layout == "1x3":
        return 1, 3, (12, 4)

    if requested_layout == "2x2":
        return 2, 2, (10, 8)

    # auto
    if n_panels <= 1:
        return 1, 1, (5.5, 4.5)
    if n_panels == 2:
        return 1, 2, (10, 4.5)
    if n_panels == 3:
        return 1, 3, (12, 4)
    return 2, 2, (10, 8)


def make_composite_figure(exp_path, selected_paths, figure_name="Figure1", layout="auto", make_pptx=True):
    exp_path = Path(exp_path)
    out_dir = exp_path / "PublicationFigures" / figure_name
    out_dir.mkdir(parents=True, exist_ok=True)

    selected_paths = [Path(p) for p in selected_paths]
    selected_paths = selected_paths[:4]

    if len(selected_paths) == 0:
        raise ValueError("No figure files were selected.")

    nrows, ncols, figsize = get_layout(len(selected_paths), requested_layout=layout)

    fig, axes = plt.subplots(nrows, ncols, figsize=figsize)

    try:
        axes_list = list(axes.ravel())
    except Exception:
        axes_list = [axes]

    letters = "ABCDEFGHIJKLMNOPQRSTUVWXYZ"
    panel_labels = []

    for i, ax in enumerate(axes_list):
        ax.axis("off")

        if i < len(selected_paths):
            img = mpimg.imread(selected_paths[i])
            ax.imshow(img)
            ax.text(
                0.01, 0.98, letters[i],
                transform=ax.transAxes,
                va="top",
                ha="left",
                fontsize=16,
                fontweight="bold"
            )
            panel_labels.append((letters[i], selected_paths[i].name))

    fig.tight_layout()

    png_path = out_dir / f"{figure_name}.png"
    svg_path = out_dir / f"{figure_name}.svg"

    fig.savefig(png_path, dpi=300, bbox_inches="tight")
    fig.savefig(svg_path, bbox_inches="tight")
    plt.close(fig)

    used_df = pd.DataFrame({
        "Panel": [p[0] for p in panel_labels],
        "FigureFile": [p.name for p in selected_paths],
        "SourcePath": [str(p) for p in selected_paths],
    })

    used_df.to_excel(out_dir / f"{figure_name}_used_files.xlsx", index=False)
    used_df.to_csv(out_dir / f"{figure_name}_used_files.csv", index=False)

    caption = generate_caption(panel_labels)
    caption_path = out_dir / f"{figure_name}_caption.txt"
    caption_path.write_text(caption, encoding="utf-8")

    pptx_path = None
    pptx_message = "PPTX was not generated."

    if make_pptx and PPTX_AVAILABLE:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(png_path), Inches(0.6), Inches(0.5), width=Inches(12.1))

        tx = slide.shapes.add_textbox(Inches(0.6), Inches(6.75), Inches(12.0), Inches(0.5))
        tf = tx.text_frame
        tf.text = figure_name
        tf.paragraphs[0].font.size = Pt(14)

        pptx_path = out_dir / f"{figure_name}.pptx"
        prs.save(pptx_path)
        pptx_message = "PPTX generated."
    elif make_pptx and not PPTX_AVAILABLE:
        pptx_message = "PPTX skipped: python-pptx is not installed."

    return {
        "output_dir": str(out_dir),
        "png": str(png_path),
        "svg": str(svg_path),
        "pptx": str(pptx_path) if pptx_path else None,
        "pptx_available": PPTX_AVAILABLE,
        "pptx_message": pptx_message,
        "caption": caption,
        "caption_path": str(caption_path),
        "used_files": used_df,
    }
